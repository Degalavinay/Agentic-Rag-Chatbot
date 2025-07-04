import os
from typing import List
from PyPDF2 import PdfReader
from docx import Document
import pandas as pd
from pptx import Presentation

class DocumentParser:
    @staticmethod
    def parse(file_path: str) -> str:
        """Parse various document formats into text"""
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == '.pdf':
            return DocumentParser._parse_pdf(file_path)
        elif ext == '.docx':
            return DocumentParser._parse_docx(file_path)
        elif ext == '.csv':
            return DocumentParser._parse_csv(file_path)
        elif ext == '.pptx':
            return DocumentParser._parse_pptx(file_path)
        elif ext in ['.txt', '.md']:
            return DocumentParser._parse_txt(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")

    @staticmethod
    def _parse_pdf(file_path: str) -> str:
        with open(file_path, 'rb') as f:
            reader = PdfReader(f)
            return "\n".join(page.extract_text() for page in reader.pages)

    @staticmethod
    def _parse_docx(file_path: str) -> str:
        doc = Document(file_path)
        return "\n".join(para.text for para in doc.paragraphs)

    @staticmethod
    def _parse_csv(file_path: str) -> str:
        df = pd.read_csv(file_path)
        return df.to_string()

    @staticmethod
    def _parse_pptx(file_path: str) -> str:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    @staticmethod
    def _parse_txt(file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
